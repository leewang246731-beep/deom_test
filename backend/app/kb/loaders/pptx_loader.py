"""PPTX Loader — using python-pptx."""
from app.kb.loaders.base_loader import BaseDocumentLoader, LoadedDocument


class PptxLoader(BaseDocumentLoader):
    @staticmethod
    def supported_extensions() -> list[str]:
        return [".pptx", ".ppt"]

    def load(self, file_path: str) -> LoadedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(file_path)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if slide_parts:
                slides_text.append(f"[Slide {i+1}]\n" + "\n".join(slide_parts))

        return LoadedDocument(
            text="\n\n".join(slides_text),
            metadata={
                "source": file_path,
                "format": "pptx",
                "slide_count": len(prs.slides),
            },
        )
